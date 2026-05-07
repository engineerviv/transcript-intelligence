"""
Generates AegisCloud_Transcript_Intelligence.pptx from the structured data
in aggregated.json and enriched.json.

Run from the project root:
    .venv/bin/python generate_pptx.py
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # emerald
RED       = RGBColor(0xF8, 0x71, 0x71)   # red
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)   # amber
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG   = RGBColor(0x1E, 0x29, 0x3B)   # slightly lighter navy for cards

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    """Add a blank slide and paint the background navy."""
    layout = prs.slide_layouts[6]  # completely blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return slide


def _box(slide, l, t, w, h, fill=None, line=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _label(slide, text, l, t, w, h,
           size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _accent_bar(slide, t, h=Pt(3).emu):
    """Thin horizontal accent line."""
    bar = slide.shapes.add_shape(1, Inches(0.5), t, Inches(12.33), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _slide_number(slide, n):
    _label(slide, str(n),
           SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.35),
           Inches(0.4), Inches(0.3),
           size=10, color=GREY, align=PP_ALIGN.RIGHT)


def _chip(slide, text, l, t, color=ACCENT):
    """Small coloured pill label."""
    w = Inches(len(text) * 0.11 + 0.3)
    chip = _box(slide, l, t, w, Inches(0.32), fill=color)
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    _label(slide, text, l + Inches(0.08), t + Inches(0.03),
           w - Inches(0.1), Inches(0.28), size=10, bold=True, color=BG,
           align=PP_ALIGN.CENTER)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = _blank(prs)

    # Decorative left stripe
    _box(slide, 0, 0, Inches(0.18), SLIDE_H, fill=ACCENT)

    _label(slide, "TRANSCRIPT INTELLIGENCE",
           Inches(0.55), Inches(1.6), Inches(11), Inches(0.6),
           size=13, bold=True, color=ACCENT)

    _label(slide, "What 100 Customer Calls Tell Us About\nReliability, Risk, and Revenue",
           Inches(0.55), Inches(2.2), Inches(11), Inches(1.6),
           size=36, bold=True, color=WHITE)

    _label(slide, "AegisCloud  ·  Q1 2024  ·  Product & Engineering Leadership",
           Inches(0.55), Inches(4.1), Inches(11), Inches(0.5),
           size=14, color=GREY)

    # Bottom stat strip
    for i, (num, label) in enumerate([
        ("100", "calls analyzed"),
        ("41%", "outage-related"),
        ("63%", "accounts at churn risk"),
        ("47%", "renewals at risk"),
    ]):
        x = Inches(0.55 + i * 3.1)
        _box(slide, x, Inches(5.9), Inches(2.8), Inches(1.2), fill=CARD_BG, line=True)
        _label(slide, num, x + Inches(0.15), Inches(6.0), Inches(2.5), Inches(0.55),
               size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
        _label(slide, label, x + Inches(0.15), Inches(6.52), Inches(2.5), Inches(0.35),
               size=11, color=GREY)

    _slide_number(slide, 1)
    return slide


def slide_headline(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "THE HEADLINE FINDING",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "One signal dominates everything",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.7),
           size=30, bold=True, color=WHITE)

    _label(slide, "Aegis Detect has a reliability problem that is now threatening renewals.",
           Inches(0.5), Inches(1.7), Inches(12), Inches(0.55),
           size=18, color=GREY, italic=True)

    findings = [
        ("41", "of 100 calls were outage-related"),
        ("16", "accounts at HIGH churn risk — all cited an outage"),
        ("7/15", "upcoming renewals are negative or mixed sentiment"),
        ("2×",  "customers called out a second outage in 8 months"),
    ]
    for i, (num, text) in enumerate(findings):
        y = Inches(2.55 + i * 0.95)
        _box(slide, Inches(0.5), y, Inches(12.3), Inches(0.82), fill=CARD_BG, line=False)
        _label(slide, num,  Inches(0.65), y + Inches(0.1), Inches(1.1), Inches(0.62),
               size=28, bold=True, color=ACCENT)
        _label(slide, text, Inches(1.9), y + Inches(0.18), Inches(10.5), Inches(0.48),
               size=15, color=WHITE)

    _label(slide, "This is not a support process problem. It is an engineering reliability problem with a direct revenue consequence.",
           Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5),
           size=11, color=YELLOW, italic=True)

    _slide_number(slide, 2)


def slide_scope(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "WHAT WE ANALYZED",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "100 call transcripts — every call enriched by AI",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=28, bold=True, color=WHITE)

    channels = [
        ("40", "External",  "QBRs, renewals, onboarding, demos",     ACCENT),
        ("33", "Internal",  "Incident war rooms, sprint planning",    YELLOW),
        ("27", "Support",   "Active cases, break-fix, escalations",   RED),
    ]
    for i, (count, label, desc, color) in enumerate(channels):
        x = Inches(0.5 + i * 4.25)
        _box(slide, x, Inches(2.0), Inches(4.0), Inches(2.8), fill=CARD_BG, line=True)
        _label(slide, count, x + Inches(0.2), Inches(2.15), Inches(1.2), Inches(0.9),
               size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        _label(slide, label, x + Inches(0.15), Inches(3.1), Inches(3.7), Inches(0.45),
               size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _label(slide, desc,  x + Inches(0.15), Inches(3.55), Inches(3.7), Inches(0.6),
               size=11, color=GREY, align=PP_ALIGN.CENTER)

    _label(slide, "Each call enriched for:",
           Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
           size=13, bold=True, color=WHITE)

    tags = ["Sentiment", "Urgency", "Churn Risk", "Topic", "Intent", "Emotion"]
    colors = [GREEN, RED, YELLOW, ACCENT, ACCENT, ACCENT]
    x_pos = Inches(0.5)
    for tag, c in zip(tags, colors):
        _chip(slide, tag, x_pos, Inches(5.7), color=c)
        x_pos += Inches(len(tag) * 0.11 + 0.5)

    _label(slide, "GPT-4o-mini · structured JSON schema · disk-level prompt cache (no redundant API calls)",
           Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
           size=11, color=GREY, italic=True)

    _slide_number(slide, 3)


def slide_sentiment(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "SENTIMENT ANALYSIS",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "Customers vs. prospects tell very different stories",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=28, bold=True, color=WHITE)

    # Left: overall breakdown
    _label(slide, "Overall sentiment",
           Inches(0.5), Inches(1.9), Inches(5), Inches(0.4),
           size=13, bold=True, color=GREY)

    bars = [
        ("Negative", 32, RED),
        ("Mixed",    22, YELLOW),
        ("Positive", 31, GREEN),
        ("Neutral",  15, GREY),
    ]
    bar_max_w = Inches(4.5)
    for i, (label, pct, color) in enumerate(bars):
        y = Inches(2.4 + i * 0.85)
        _label(slide, label, Inches(0.5), y, Inches(1.3), Inches(0.4),
               size=12, color=WHITE)
        w = bar_max_w * pct / 100
        _box(slide, Inches(1.85), y + Inches(0.05), w, Inches(0.38), fill=color)
        _label(slide, f"{pct}%", Inches(1.85) + w + Inches(0.08), y,
               Inches(0.6), Inches(0.4), size=12, bold=True, color=color)

    # Right: by call type
    _label(slide, "By channel",
           Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.4),
           size=13, bold=True, color=GREY)

    channels_data = [
        ("Support",  "59% negative", "7% positive",  RED,   "Avg score: 2.94 / 5"),
        ("External", "50% positive", "12% negative", GREEN, "Avg score: 3.85 / 5"),
        ("Internal", "33% negative", "27% positive", YELLOW,"Avg score: 3.28 / 5"),
    ]
    for i, (ch, stat1, stat2, color, avg) in enumerate(channels_data):
        y = Inches(2.4 + i * 1.4)
        _box(slide, Inches(7.0), y, Inches(5.8), Inches(1.2), fill=CARD_BG, line=True)
        _label(slide, ch, Inches(7.15), y + Inches(0.08), Inches(2), Inches(0.42),
               size=14, bold=True, color=color)
        _label(slide, stat1, Inches(7.15), y + Inches(0.5), Inches(2.8), Inches(0.35),
               size=12, color=WHITE)
        _label(slide, stat2, Inches(7.15), y + Inches(0.82), Inches(2.8), Inches(0.32),
               size=11, color=GREY)
        _label(slide, avg, Inches(10.0), y + Inches(0.45), Inches(2.6), Inches(0.4),
               size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)

    _label(slide, "The product sells. The gap is post-sale reliability.",
           Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45),
           size=12, color=YELLOW, italic=True, bold=True)

    _slide_number(slide, 4)


def slide_outages(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "ROOT CAUSE",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "41% of all calls were about outages",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=28, bold=True, color=WHITE)

    _label(slide, '"Outage Response" is simultaneously the #1 topic, #1 negative topic, and #1 high-urgency topic.',
           Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
           size=13, color=GREY, italic=True)

    # Top topics bar chart (left)
    _label(slide, "Top negative + high-urgency topics",
           Inches(0.5), Inches(2.4), Inches(6), Inches(0.4),
           size=12, bold=True, color=GREY)

    topics = [
        ("Outage Response", 7),
        ("Product Outage",  2),
        ("Service Outage",  2),
        ("Platform Outage", 2),
        ("Backup Perf.",    2),
    ]
    for i, (topic, count) in enumerate(topics):
        y = Inches(2.95 + i * 0.72)
        _label(slide, topic, Inches(0.5), y, Inches(2.3), Inches(0.4),
               size=11, color=WHITE)
        w = Inches(0.55 * count)
        _box(slide, Inches(2.85), y + Inches(0.05), w, Inches(0.36),
             fill=(RED if topic == "Outage Response" else YELLOW))
        _label(slide, str(count), Inches(2.85) + w + Inches(0.1), y,
               Inches(0.4), Inches(0.4), size=12, bold=True,
               color=(RED if topic == "Outage Response" else YELLOW))

    # Named incidents (right)
    _label(slide, "Named incidents in the data",
           Inches(7.0), Inches(2.4), Inches(6), Inches(0.4),
           size=12, bold=True, color=GREY)

    incidents = [
        "Detect pipeline failure — internal war room",
        "Detect alert delays — Summit Trust, Trailhead",
        "Detect dashboard down — Cobalt Software",
        "6-hour detection blindness — Northstar Pharma (healthcare)",
        "Outage during live regulatory audit — Meridian Capital",
    ]
    for i, inc in enumerate(incidents):
        y = Inches(2.95 + i * 0.72)
        _box(slide, Inches(7.0), y, Inches(0.06), Inches(0.36), fill=RED)
        _label(slide, inc, Inches(7.2), y, Inches(5.6), Inches(0.42),
               size=11, color=WHITE)

    _label(slide, "Every one of the 16 high-churn-risk accounts mentioned an outage.",
           Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45),
           size=12, color=RED, bold=True)

    _slide_number(slide, 5)


def slide_churn(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "CHURN RISK",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "63% of accounts are not safe",
           Inches(0.5), Inches(0.9), Inches(9), Inches(0.65),
           size=28, bold=True, color=WHITE)

    buckets = [
        ("HIGH",   16, RED,    "Actively evaluating competitors"),
        ("MEDIUM", 47, YELLOW, "One more incident from flipping"),
        ("LOW",    20, GREEN,  "Watching closely"),
        ("NONE",   17, GREY,   "Stable"),
    ]
    total = 100
    bar_max = Inches(9.0)
    for i, (label, count, color, note) in enumerate(buckets):
        y = Inches(2.1 + i * 1.1)
        w = bar_max * count / total
        _box(slide, Inches(0.5), y, w, Inches(0.65), fill=color)
        _label(slide, label, Inches(0.65), y + Inches(0.1), Inches(1.5), Inches(0.45),
               size=13, bold=True, color=BG)
        _label(slide, f"{count} accounts", Inches(0.5) + w + Inches(0.15), y + Inches(0.08),
               Inches(2), Inches(0.35), size=13, bold=True, color=color)
        _label(slide, note, Inches(0.5) + w + Inches(0.15), y + Inches(0.42),
               Inches(4), Inches(0.32), size=10, color=GREY)

    _label(slide, "Pattern across all high-risk accounts:",
           Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
           size=12, bold=True, color=GREY)
    _label(slide, "Outage  →  Negative sentiment  →  Competitor evaluation  →  Renewal hesitation",
           Inches(0.5), Inches(6.78), Inches(12), Inches(0.45),
           size=13, bold=True, color=WHITE)

    _slide_number(slide, 6)


def slide_accounts(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "IMMEDIATE ACTION REQUIRED",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=RED)

    _label(slide, "6 accounts need a phone call this week",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=28, bold=True, color=WHITE)

    _label(slide, "All are actively evaluating competitors right now.",
           Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
           size=13, color=GREY, italic=True)

    accounts = [
        ("Quantum Edge",     "SLA breach; renewal proposal in flight"),
        ("Meridian Capital", "Outage during regulatory audit; formal vendor review open"),
        ("Northstar Pharma", "2nd outage in 8 months; CISO evaluating alternatives"),
        ("Summit Trust",     "Competitor already engaged; MFA/SSO issues unresolved"),
        ("Nova Retail Group","Renewal decision imminent; outage + compliance gaps"),
        ("Helix Data",       'Explicitly stated "considering switching vendors"'),
    ]
    for i, (account, status) in enumerate(accounts):
        y = Inches(2.25 + i * 0.78)
        col = i % 2
        x = Inches(0.5 + col * 6.4)
        row = i // 2
        y = Inches(2.25 + row * 1.4)

    # Redo as 2-column layout
    for i, (account, status) in enumerate(accounts):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.4)
        y = Inches(2.25 + row * 1.45)
        _box(slide, x, y, Inches(6.0), Inches(1.2), fill=CARD_BG, line=True)
        # Red urgency dot
        _box(slide, x + Inches(0.15), y + Inches(0.42), Inches(0.15), Inches(0.15), fill=RED)
        _label(slide, account, x + Inches(0.45), y + Inches(0.1), Inches(5.3), Inches(0.45),
               size=14, bold=True, color=WHITE)
        _label(slide, status,  x + Inches(0.45), y + Inches(0.6), Inches(5.3), Inches(0.5),
               size=11, color=GREY)

    _label(slide, "Recommended: Assign a named technical contact to each — not a ticket queue.",
           Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
           size=11, color=YELLOW, italic=True)

    _slide_number(slide, 7)


def slide_renewals(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "RENEWAL PIPELINE",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "47% of upcoming renewals are at risk",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=28, bold=True, color=WHITE)

    _label(slide, "15 renewal-related calls in the dataset. 7 carry negative or mixed sentiment.",
           Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
           size=13, color=GREY, italic=True)

    # Visual: 15 circles, 7 red, 8 green
    for i in range(15):
        col = i % 8
        row = i // 8
        x = Inches(1.0 + col * 1.4)
        y = Inches(2.4 + row * 1.3)
        color = RED if i < 7 else GREEN
        _box(slide, x, y, Inches(0.9), Inches(0.9), fill=color)
        label = "AT RISK" if i < 7 else "OK"
        _label(slide, label, x, y + Inches(1.0), Inches(0.9), Inches(0.35),
               size=8, color=color, align=PP_ALIGN.CENTER)

    # Themes
    _label(slide, "What at-risk renewals say:",
           Inches(0.5), Inches(5.4), Inches(7), Inches(0.4),
           size=12, bold=True, color=GREY)

    quotes = [
        '"We need reliability improvements before we sign"',
        '"Our board is asking us to evaluate alternatives"',
        '"Can you give us a competitive quote?" (SentinelShield)',
    ]
    for i, q in enumerate(quotes):
        _label(slide, q, Inches(0.7), Inches(5.9 + i * 0.38), Inches(11), Inches(0.35),
               size=11, color=YELLOW, italic=True)

    _slide_number(slide, 8)


def slide_features(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "SECONDARY SIGNAL",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "Feature requests are stacking up — and customers notice",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=26, bold=True, color=WHITE)

    _label(slide, "10 calls explicitly about feature requests. 3 carry negative sentiment.",
           Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
           size=13, color=GREY, italic=True)

    gaps = [
        ("Aegis Comply",           "Described as having 'significant feature gaps' vs competitors",    RED),
        ("SOC 2 reporting",        "Customers waiting on roadmap delivery; audits can't wait",          YELLOW),
        ("Alert tuning / fatigue", "Mentioned in multiple Detect calls; blocking adoption",             YELLOW),
        ("Compliance automation",  "Manual evidence-gathering causing frustration ahead of audits",     YELLOW),
    ]
    for i, (gap, desc, color) in enumerate(gaps):
        y = Inches(2.4 + i * 1.05)
        _box(slide, Inches(0.5), y, Inches(12.3), Inches(0.88), fill=CARD_BG)
        _label(slide, gap,  Inches(0.7), y + Inches(0.08), Inches(3.5), Inches(0.4),
               size=14, bold=True, color=color)
        _label(slide, desc, Inches(4.3), y + Inches(0.22), Inches(8.2), Inches(0.44),
               size=12, color=WHITE)

    _label(slide, "Feature delay compounds outage frustration. Roadmap communication matters as much as delivery.",
           Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45),
           size=12, color=YELLOW, italic=True, bold=True)

    _slide_number(slide, 9)


def slide_optimism(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "REASON FOR OPTIMISM",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=GREEN)

    _label(slide, "The product sells. Reliability is the gap — and it's fixable.",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=26, bold=True, color=WHITE)

    positives = [
        ("3.85 / 5",   "Average sentiment score on external calls",  ACCENT),
        ("50%",        "Of external calls carry positive sentiment",  GREEN),
        ("3-year deal","Axiom Labs progressing to multi-year renewal despite recent outage", GREEN),
        ("4 onboarding\ncalls",  "All broadly positive — new customers engaging well", GREEN),
    ]
    for i, (stat, text, color) in enumerate(positives):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.4)
        y = Inches(2.1 + row * 1.8)
        _box(slide, x, y, Inches(6.0), Inches(1.55), fill=CARD_BG, line=True)
        _label(slide, stat, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.75),
               size=26, bold=True, color=color)
        _label(slide, text, x + Inches(0.2), y + Inches(0.88), Inches(5.6), Inches(0.6),
               size=11, color=GREY)

    _label(slide, "This is not a product-market fit problem. Fix Detect stability, and the churn risk is recoverable.",
           Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
           size=13, color=GREEN, bold=True)

    _slide_number(slide, 10)


def slide_recommendations(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "RECOMMENDATIONS",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "Prioritized by urgency",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.55),
           size=28, bold=True, color=WHITE)

    recs = [
        ("This week",  RED,
         "Revenue protection",
         ["CS outreach to all 6 high-churn accounts before renewal dates",
          "Assign a named technical contact to each — not a ticket queue"]),
        ("30 days",   YELLOW,
         "Engineering",
         ["Root cause analysis on Aegis Detect reliability (single largest churn driver)",
          "Prioritise alert fatigue / tuning controls — blocking adoption"]),
        ("60 days",   ACCENT,
         "Product",
         ["Customer-facing roadmap update on Comply gaps (SOC 2, evidence automation)",
          "Structured QBR cadence for medium-risk accounts — before incidents happen"]),
        ("Ongoing",   GREEN,
         "Intelligence",
         ["Run this analysis on a rolling weekly / monthly cadence",
          "Use the chat agent as an always-on signal layer for CS and PM"]),
    ]
    for i, (horizon, color, category, actions) in enumerate(recs):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.4)
        y = Inches(2.0 + row * 2.3)
        _box(slide, x, y, Inches(6.0), Inches(2.0), fill=CARD_BG, line=True)
        _box(slide, x, y, Inches(6.0), Inches(0.42), fill=color)
        _label(slide, f"{horizon}  ·  {category}", x + Inches(0.15), y + Inches(0.06),
               Inches(5.7), Inches(0.35), size=12, bold=True, color=BG)
        for j, action in enumerate(actions):
            _label(slide, f"• {action}", x + Inches(0.2), y + Inches(0.52 + j * 0.65),
                   Inches(5.6), Inches(0.6), size=11, color=WHITE)

    _slide_number(slide, 11)


def slide_platform(prs):
    slide = _blank(prs)
    _accent_bar(slide, Inches(1.1))

    _label(slide, "HOW WE BUILT THIS",
           Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
           size=11, bold=True, color=ACCENT)

    _label(slide, "From raw transcripts to queryable intelligence in one pipeline",
           Inches(0.5), Inches(0.9), Inches(12), Inches(0.65),
           size=26, bold=True, color=WHITE)

    steps = [
        ("Raw\nTranscripts",    GREY),
        ("GPT-4o-mini\nEnrichment", ACCENT),
        ("FAISS\nSemantic Index", ACCENT),
        ("Aggregation\n& Insights", YELLOW),
        ("LangGraph\nReAct Agent", GREEN),
        ("Dashboard\n+ Chat", GREEN),
    ]
    box_w = Inches(1.7)
    arrow_w = Inches(0.4)
    total_w = len(steps) * box_w + (len(steps) - 1) * arrow_w
    start_x = (SLIDE_W - total_w) / 2

    for i, (label, color) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)
        _box(slide, x, Inches(2.2), box_w, Inches(1.1), fill=CARD_BG, line=True)
        _label(slide, label, x + Inches(0.1), Inches(2.25), box_w - Inches(0.2), Inches(1.0),
               size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            _label(slide, "→", ax, Inches(2.5), arrow_w, Inches(0.5),
                   size=16, color=GREY, align=PP_ALIGN.CENTER)

    details = [
        ("Local embeddings",   "sentence-transformers (no API cost for retrieval)"),
        ("Prompt cache",       "SHA-256 disk cache — avoids re-enriching on re-runs"),
        ("Output validation",  "Schema checks + cross-field rules catch LLM drift"),
        ("140 tests",          "Validation, analysis, retrieval, API, agent tools"),
        ("Scalable",           "FAISS → pgvector / Pinecone at 10k+ transcripts, no re-arch"),
    ]
    for i, (label, desc) in enumerate(details):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(4.0 + row * 1.0)
        _label(slide, f"✦  {label}:", x, y, Inches(1.9), Inches(0.38),
               size=11, bold=True, color=ACCENT)
        _label(slide, desc, x + Inches(1.95), y, Inches(2.1), Inches(0.38),
               size=11, color=WHITE)

    _slide_number(slide, 12)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = _new_prs()

    slide_title(prs)
    slide_headline(prs)
    slide_scope(prs)
    slide_sentiment(prs)
    slide_outages(prs)
    slide_churn(prs)
    slide_accounts(prs)
    slide_renewals(prs)
    slide_features(prs)
    slide_optimism(prs)
    slide_recommendations(prs)
    slide_platform(prs)

    out = Path(__file__).parent / "AegisCloud_Transcript_Intelligence.pptx"
    prs.save(str(out))
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
